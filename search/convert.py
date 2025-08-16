from pptx import Presentation
from docx import Document
from PyPDF2 import PdfReader
import pandas as pd
import re


class Converter:
    
    @staticmethod
    def remove_spaces(text):
        return re.sub(' +', ' ', text)


    @staticmethod
    def remove_nan(text):
        return text.replace('NaN', ' ')
    

    @staticmethod
    def remove_tab(text):
        return text.replace('\t', ' ')
    

    @staticmethod
    def remove_n(text):
        return text.replace('\n', ' ')
    
    
    @staticmethod
    def __get_docx_text(docx_path):
        doc = Document(docx_path)
        full_text = []

        for para in doc.paragraphs:
            full_text.append(para.text)

        return "\n".join(full_text)


    @staticmethod
    def __get_pptx_text(presentation_path):
        prs = Presentation(presentation_path)
        text_runs = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)

        return "\n".join(text_runs)
    
    
    @staticmethod
    def __get_pdf_text(pdf_path):
        reader = PdfReader(pdf_path)
        text = ""

        for page in reader.pages:
            text += page.extract_text()

        return text
    
    
    @classmethod
    def __read_excel_file(self, excel_path):
        df = pd.read_excel(excel_path, header=None)
        text = df.to_string(index=False, header=False)
        return text
    
    
    @classmethod
    def convert(self, path):
        f = {
            'pptx' : self.__get_pptx_text,
            'docx' : self.__get_docx_text,
            'pdf'  : self.__get_pdf_text,
        }
        
        extension = path.rsplit('.', 1)[1].lower()
        
        try:
            result = f[extension](path)
        except:
            result = open(path, 'r').read()
        
        return self.remove_n(self.remove_nan(self.remove_tab(self.remove_spaces(result))))